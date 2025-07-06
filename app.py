from flask import Flask, request, render_template, jsonify, flash, redirect, url_for
import os
import PyPDF2
from pptx import Presentation
import requests
from werkzeug.utils import secure_filename
import re
from youtube_search import YoutubeSearch
import json
from config import config

app = Flask(__name__)

# Load configuration
config_name = os.environ.get('FLASK_ENV') or 'default'
app.config.from_object(config[config_name])

# Ensure upload directory exists
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

ALLOWED_EXTENSIONS = app.config['ALLOWED_EXTENSIONS']

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_pdf(file_path):
    """Extract text from PDF file"""
    text = ""
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text()
    except Exception as e:
        print(f"Error reading PDF: {e}")
    return text

def extract_text_from_ppt(file_path):
    """Extract text from PowerPoint file"""
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error reading PPT: {e}")
    return text

def extract_topics_from_text(text):
    """Extract main topics from text using simple keyword extraction"""
    # Remove extra whitespace and convert to lowercase
    text = re.sub(r'\s+', ' ', text.lower())
    
    # Simple topic extraction - looking for common academic patterns
    topics = []
    
    # Get configuration values
    min_length = app.config.get('MIN_TOPIC_LENGTH', 10)
    max_length = app.config.get('MAX_TOPIC_LENGTH', 100)
    max_topics = app.config.get('MAX_TOPICS_TO_EXTRACT', 10)
    
    # Look for chapter/section headers using configured patterns
    chapter_patterns = app.config.get('CHAPTER_PATTERNS', [
        r'chapter\s+\d+[:\-\s]+([^\n\.]+)',
        r'section\s+\d+[:\-\s]+([^\n\.]+)',
        r'unit\s+\d+[:\-\s]+([^\n\.]+)',
        r'topic\s+\d+[:\-\s]+([^\n\.]+)',
        r'lecture\s+\d+[:\-\s]+([^\n\.]+)'
    ])
    
    for pattern in chapter_patterns:
        matches = re.findall(pattern, text)
        topics.extend([match.strip() for match in matches])
    
    # Look for bullet points or numbered lists using configured patterns
    bullet_patterns = app.config.get('BULLET_PATTERNS', [
        r'•\s*([^\n•]+)',
        r'\d+\.\s*([^\n\d]+)',
        r'-\s*([^\n-]+)'
    ])
    
    for pattern in bullet_patterns:
        matches = re.findall(pattern, text)
        # Filter out matches that don't meet length requirements
        topics.extend([match.strip() for match in matches if min_length <= len(match.strip()) <= max_length])
    
    # Extract sentences with academic keywords
    academic_keywords = app.config.get('ACADEMIC_KEYWORDS', [
        'introduction', 'concept', 'theory', 'principle', 'method', 'approach',
        'algorithm', 'technique', 'framework', 'model', 'system', 'process',
        'analysis', 'implementation', 'application', 'example', 'case study'
    ])
    
    sentences = text.split('.')
    for sentence in sentences:
        sentence = sentence.strip()
        if any(keyword in sentence.lower() for keyword in academic_keywords):
            if min_length <= len(sentence) <= max_length:
                topics.append(sentence)
    
    # Remove duplicates and return unique topics
    unique_topics = list(set(topics))
    return unique_topics[:max_topics]

def search_youtube_videos(query, max_results=5):
    """Search YouTube for videos related to the query"""
    try:
        # Clean the query
        query = re.sub(r'[^\w\s]', '', query)
        query = query.strip()
        
        # Use youtube-search library
        results = YoutubeSearch(query, max_results=max_results).to_dict()
        
        videos = []
        for result in results:
            video_info = {
                'title': result['title'],
                'url': 'https://youtube.com' + result['url_suffix'],
                'duration': result['duration'],
                'views': result['views'],
                'channel': result['channel'],
                'thumbnail': result['thumbnails'][0] if result['thumbnails'] else None
            }
            videos.append(video_info)
        
        return videos
    except Exception as e:
        print(f"Error searching YouTube: {e}")
        return []

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        flash('No file selected')
        return redirect(request.url)
    
    file = request.files['file']
    if file.filename == '':
        flash('No file selected')
        return redirect(request.url)
    
    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)
        
        # Extract text based on file type
        if filename.lower().endswith('.pdf'):
            text = extract_text_from_pdf(file_path)
        elif filename.lower().endswith(('.ppt', '.pptx')):
            text = extract_text_from_ppt(file_path)
        else:
            flash('Unsupported file type')
            return redirect(url_for('index'))
        
        # Extract topics
        topics = extract_topics_from_text(text)
        
        # Search for YouTube videos for each topic
        video_results = {}
        max_videos = app.config.get('MAX_VIDEOS_PER_TOPIC', 5)
        for topic in topics:
            if topic.strip():  # Only search for non-empty topics
                videos = search_youtube_videos(topic, max_videos)
                if videos:
                    video_results[topic] = videos
        
        # Clean up uploaded file
        os.remove(file_path)
        
        return render_template('results.html', 
                             filename=filename, 
                             topics=topics,
                             video_results=video_results)
    
    flash('Invalid file type. Please upload PDF, PPT, or PPTX files.')
    return redirect(url_for('index'))

@app.route('/api/search', methods=['POST'])
def api_search():
    """API endpoint for searching videos by topic"""
    data = request.get_json()
    query = data.get('query', '')
    max_results = data.get('max_results', 5)
    
    videos = search_youtube_videos(query, max_results)
    return jsonify({'videos': videos})

if __name__ == '__main__':
    app.run(debug=True)
